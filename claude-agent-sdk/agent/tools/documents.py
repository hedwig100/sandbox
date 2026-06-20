"""Readers for unstructured Office documents (Excel / PowerPoint).

These tools turn binary `.xlsx` / `.pptx` files into plain text the model can
read. They use the JSON-Schema form of the input spec (rather than the simple
``{"name": type}`` map) so optional parameters like ``sheet`` are expressible.

`path` defaults to the bundled sample file from `agent.config`, so the agent can
call the tool with no arguments to inspect the demo data.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from claude_agent_sdk import create_sdk_mcp_server, tool

from agent.config import PPTX_PATH, XLSX_PATH

MAX_ROWS = 100


def _err(text: str) -> dict[str, Any]:
    return {"content": [{"type": "text", "text": text}], "is_error": True}


@tool(
    "read_excel",
    "Read an Excel (.xlsx) workbook as text. Returns the named sheet as a "
    "markdown table, or a list of sheet names plus the first sheet when no "
    "sheet is given.",
    {
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Path to the .xlsx file. Defaults to the bundled sample.",
            },
            "sheet": {
                "type": "string",
                "description": "Sheet name to read. Omit to read the active sheet.",
            },
        },
    },
)
async def read_excel(args: dict[str, Any]) -> dict[str, Any]:
    from openpyxl import load_workbook

    path = Path(args.get("path") or XLSX_PATH)
    if not path.exists():
        return _err(f"Excel file not found: {path}. Run `poetry run seed` first.")

    wb = load_workbook(path, read_only=True, data_only=True)
    sheet = args.get("sheet")
    if sheet and sheet not in wb.sheetnames:
        return _err(f"Unknown sheet {sheet!r}. Available: {wb.sheetnames}")
    ws = wb[sheet] if sheet else wb.active

    rows: list[list[str]] = []
    for row in ws.iter_rows(max_row=MAX_ROWS, values_only=True):
        rows.append(["" if v is None else str(v) for v in row])

    header = f"Workbook sheets: {wb.sheetnames}\nShowing sheet: {ws.title}\n\n"
    if not rows:
        return {"content": [{"type": "text", "text": header + "(empty sheet)"}]}

    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    lines = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join("---" for _ in range(width)) + " |"]
    lines += ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return {"content": [{"type": "text", "text": header + "\n".join(lines)}]}


@tool(
    "read_powerpoint",
    "Read a PowerPoint (.pptx) deck as text. Returns the text of every slide "
    "(titles, bullets, tables, and speaker notes).",
    {
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Path to the .pptx file. Defaults to the bundled sample.",
            },
        },
    },
)
async def read_powerpoint(args: dict[str, Any]) -> dict[str, Any]:
    from pptx import Presentation

    path = Path(args.get("path") or PPTX_PATH)
    if not path.exists():
        return _err(f"PowerPoint file not found: {path}. Run `poetry run seed` first.")

    prs = Presentation(str(path))
    blocks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    lines.append("| " + " | ".join(c.text for c in row.cells) + " |")
            elif shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"_Notes: {notes}_")
        blocks.append("\n".join(lines))

    body = "\n\n".join(blocks) if blocks else "(deck has no slides)"
    return {"content": [{"type": "text", "text": body}]}


DOCUMENTS_SERVER = create_sdk_mcp_server(
    name="documents",
    version="0.1.0",
    tools=[read_excel, read_powerpoint],
)

DOCUMENTS_TOOL_NAMES = [
    "mcp__documents__read_excel",
    "mcp__documents__read_powerpoint",
]
