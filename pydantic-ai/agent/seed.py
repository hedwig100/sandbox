"""Generate the sample data the agent's tools read.

Creates, under `data/` (gitignored):
  * sample.db    -> SQLite with products / customers / orders tables + rows
  * sample.xlsx  -> a 2-sheet workbook (Inventory, Sales)
  * sample.pptx  -> a short deck

Re-run any time with `uv run seed`; it overwrites the existing files.
"""

from __future__ import annotations

import sqlite3

from agent.config import DATA_DIR, DB_PATH, PPTX_PATH, XLSX_PATH

PRODUCTS = [
    (1, "Widget", "Hardware", 9.99, 3),
    (2, "Gadget", "Hardware", 19.99, 0),
    (3, "Gizmo", "Hardware", 14.50, 42),
    (4, "Notebook", "Stationery", 4.20, 5),
    (5, "Pen", "Stationery", 1.10, 250),
]
CUSTOMERS = [
    (1, "Alice Tanaka", "alice@example.com", "Tokyo"),
    (2, "Bob Suzuki", "bob@example.com", "Osaka"),
    (3, "Carol Sato", "carol@example.com", "Tokyo"),
]
ORDERS = [
    (1, 1, 3, 2, "2026-05-01"),
    (2, 2, 5, 10, "2026-05-03"),
    (3, 1, 1, 1, "2026-06-10"),
]


def _seed_db() -> None:
    DB_PATH.unlink(missing_ok=True)
    conn = sqlite3.connect(DB_PATH)
    conn.executescript(
        """
        CREATE TABLE products (
            id INTEGER PRIMARY KEY,
            name TEXT NOT NULL,
            category TEXT NOT NULL,
            price REAL NOT NULL,
            stock INTEGER NOT NULL
        );
        CREATE TABLE customers (
            id INTEGER PRIMARY KEY,
            name TEXT NOT NULL,
            email TEXT NOT NULL,
            city TEXT NOT NULL
        );
        CREATE TABLE orders (
            id INTEGER PRIMARY KEY,
            customer_id INTEGER NOT NULL REFERENCES customers(id),
            product_id INTEGER NOT NULL REFERENCES products(id),
            quantity INTEGER NOT NULL,
            ordered_at TEXT NOT NULL
        );
        """
    )
    conn.executemany("INSERT INTO products VALUES (?, ?, ?, ?, ?)", PRODUCTS)
    conn.executemany("INSERT INTO customers VALUES (?, ?, ?, ?)", CUSTOMERS)
    conn.executemany("INSERT INTO orders VALUES (?, ?, ?, ?, ?)", ORDERS)
    conn.commit()
    conn.close()
    print(f"wrote {DB_PATH}")


def _seed_xlsx() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    inv = wb.active
    inv.title = "Inventory"
    inv.append(["id", "name", "category", "price", "stock"])
    for row in PRODUCTS:
        inv.append(list(row))

    sales = wb.create_sheet("Sales")
    sales.append(["order_id", "customer", "product", "quantity", "date"])
    cust = {c[0]: c[1] for c in CUSTOMERS}
    prod = {p[0]: p[1] for p in PRODUCTS}
    for oid, cid, pid, qty, date in ORDERS:
        sales.append([oid, cust[cid], prod[pid], qty, date])

    wb.save(XLSX_PATH)
    print(f"wrote {XLSX_PATH}")


def _seed_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    title = prs.slides.add_slide(prs.slide_layouts[0])
    title.shapes.title.text = "Acme Q2 2026 Review"
    title.placeholders[1].text = "Inventory & sales snapshot"

    bullets = prs.slides.add_slide(prs.slide_layouts[1])
    bullets.shapes.title.text = "Highlights"
    body = bullets.placeholders[1].text_frame
    body.text = "Gadget is out of stock (0 units)"
    body.add_paragraph().text = "Pen is the best seller (10 units sold)"
    body.add_paragraph().text = "Tokyo is the top customer city"
    bullets.notes_slide.notes_text_frame.text = "Restock Gadget before July."

    table_slide = prs.slides.add_slide(prs.slide_layouts[5])
    table_slide.shapes.title.text = "Low stock items"
    low = [p for p in PRODUCTS if p[4] < 5]
    tbl = table_slide.shapes.add_table(
        len(low) + 1, 2, Inches(1), Inches(1.8), Inches(5), Inches(2)
    ).table
    tbl.cell(0, 0).text = "Product"
    tbl.cell(0, 1).text = "Stock"
    for r, p in enumerate(low, start=1):
        tbl.cell(r, 0).text = p[1]
        tbl.cell(r, 1).text = str(p[4])

    prs.save(PPTX_PATH)
    print(f"wrote {PPTX_PATH}")


def main() -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    _seed_db()
    _seed_xlsx()
    _seed_pptx()
    print("done.")


if __name__ == "__main__":
    main()
