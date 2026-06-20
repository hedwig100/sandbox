"""Readers for unstructured Office documents (Excel / PowerPoint).

Plain tool functions that turn binary `.xlsx` / `.pptx` files into text the model
can read. ``path`` defaults to the bundled sample from `agent.config`, so the
agent can call a reader with no arguments to inspect the demo data.
"""

from __future__ import annotations

from pathlib import Path

from agent.config import PPTX_PATH, XLSX_PATH

MAX_ROWS = 100


def read_excel(path: str = "", sheet: str = "") -> str:
    """Read an Excel (.xlsx) workbook as text.

    Returns the named sheet as a markdown table, or the active sheet (plus the
    list of sheet names) when no sheet is given.

    Args:
        path: Path to the .xlsx file. Leave empty to use the bundled sample.
        sheet: Sheet name to read. Leave empty to read the active sheet.
    """
    from openpyxl import load_workbook

    file = Path(path) if path else XLSX_PATH
    if not file.exists():
        return f"Excel file not found: {file}. Run `uv run seed` first."

    wb = load_workbook(file, read_only=True, data_only=True)
    if sheet and sheet not in wb.sheetnames:
        return f"Unknown sheet {sheet!r}. Available: {wb.sheetnames}"
    ws = wb[sheet] if sheet else wb.active

    rows: list[list[str]] = []
    for row in ws.iter_rows(max_row=MAX_ROWS, values_only=True):
        rows.append(["" if v is None else str(v) for v in row])

    header = f"Workbook sheets: {wb.sheetnames}\nShowing sheet: {ws.title}\n\n"
    if not rows:
        return header + "(empty sheet)"

    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    lines = [
        "| " + " | ".join(rows[0]) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
    ]
    lines += ["| " + " | ".join(r) + " |" for r in rows[1:]]
    return header + "\n".join(lines)


def read_powerpoint(path: str = "") -> str:
    """Read a PowerPoint (.pptx) deck as text.

    Returns the text of every slide (titles, bullets, tables, and speaker notes).

    Args:
        path: Path to the .pptx file. Leave empty to use the bundled sample.
    """
    from pptx import Presentation

    file = Path(path) if path else PPTX_PATH
    if not file.exists():
        return f"PowerPoint file not found: {file}. Run `uv run seed` first."

    prs = Presentation(str(file))
    blocks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    lines.append("| " + " | ".join(c.text for c in row.cells) + " |")
            elif shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"_Notes: {notes}_")
        blocks.append("\n".join(lines))

    return "\n\n".join(blocks) if blocks else "(deck has no slides)"


DOCUMENT_TOOLS = [read_excel, read_powerpoint]
